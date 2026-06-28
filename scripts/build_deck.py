import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

INK=RGBColor(0x14,0x25,0x3A); TEAL=RGBColor(0x0E,0x7C,0x7B); CORAL=RGBColor(0xE0,0x7A,0x5F)
PANEL=RGBColor(0xF2,0xF5,0xF7); MUTED=RGBColor(0x64,0x74,0x8B); WHITE=RGBColor(0xFF,0xFF,0xFF)
CODEBG=RGBColor(0x14,0x25,0x3A); CODEFG=RGBColor(0xE6,0xED,0xF3); CODEKEY=RGBColor(0x7F,0xD8,0xD6)
GRID=RGBColor(0xE2,0xE8,0xEE)
SANS="Segoe UI"; MONO="Consolas"

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]

import os
os.makedirs("/tmp/pptx_assets", exist_ok=True)
def _render_structures():
    from rdkit import Chem
    from rdkit.Chem import Draw
    from rdkit.Chem.Draw import rdMolDraw2D
    mols=[("Aspirin","CC(=O)Oc1ccccc1C(=O)O","BSYNRYMUTXBXSQ","NIST WebBook"),
          ("Caffeine","Cn1c(=O)c2c(ncn2C)n(C)c1=O","RYYVLZVUVIJVGH","NIST WebBook"),
          ("Ibuprofen","CC(C)Cc1ccc(cc1)C(C)C(=O)O","HEFNNWSXXWATRW","NIST WebBook"),
          ("Piperazine","C1CNCCN1","GLUUGHFHXGJENI","paper extract"),
          ("Trimesoyl chloride","O=C(Cl)c1cc(C(=O)Cl)cc(C(=O)Cl)c1","UWCPYKQBIPYOLX","paper extract"),
          ("Benzene","c1ccccc1","UHOVQNZJYSORNB","NIST WebBook")]
    meta=[]
    for name,smi,ik,src in mols:
        m=Chem.MolFromSmiles(smi); pp=f"/tmp/pptx_assets/{ik}.png"
        try:
            d=rdMolDraw2D.MolDraw2DCairo(460,340); o=d.drawOptions(); o.bondLineWidth=2; o.padding=0.12
            rdMolDraw2D.PrepareAndDrawMolecule(d,m); d.FinishDrawing(); open(pp,"wb").write(d.GetDrawingText())
        except Exception:
            Draw.MolToImage(m,size=(460,340)).save(pp)
        meta.append({"name":name,"ik":ik,"src":src,"png":pp,"smiles":Chem.MolToSmiles(m)})
    json.dump(meta, open("/tmp/pptx_assets/meta.json","w")); return meta
_render_structures()

def _mermaid_png(mmd_path, cache):
    if os.path.exists(cache): return cache
    try:
        import base64, httpx
        code=open(mmd_path).read()
        b64=base64.urlsafe_b64encode(code.encode()).decode()
        r=httpx.get(f"https://mermaid.ink/img/{b64}?type=png&bgColor=FFFFFF", timeout=60, follow_redirects=True)
        if r.status_code==200 and r.content[:4]==b"\x89PNG":
            open(cache,"wb").write(r.content); return cache
    except Exception:
        pass
    return None
SG1_PNG=_mermaid_png("/root/scinex/docs/sg1_ingest.mmd","/root/scinex/docs/sg1_ingest.png")
SG2_PNG=_mermaid_png("/root/scinex/docs/sg2_structure.mmd","/root/scinex/docs/sg2_structure.png")
SG3_PNG=_mermaid_png("/root/scinex/docs/sg3_qa.mmd","/root/scinex/docs/sg3_qa.png")

def slide(): return prs.slides.add_slide(BLANK)

def box(s,x,y,w,h,fill=None,line=None,rad=True):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rad else MSO_SHAPE.RECTANGLE,
                           Inches(x),Inches(y),Inches(w),Inches(h))
    shp.shadow.inherit=False
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(1)
    shp.text_frame.paragraphs[0].text=""
    return shp

def txt(s,x,y,w,h,paras,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP):
    tb=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h))
    tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    tf.margin_left=0; tf.margin_right=0; tf.margin_top=0; tf.margin_bottom=0
    for i,para in enumerate(paras):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align
        sb=para.get("space",None)
        if sb is not None: p.space_after=Pt(sb)
        for run in para["runs"]:
            r=p.add_run(); r.text=run[0]; f=r.font
            f.size=Pt(run[1]); f.bold=run[2]; f.color.rgb=run[3]
            f.name=run.get(4,SANS) if isinstance(run,dict) else (run[4] if len(run)>4 else SANS)
    return tb

def header(s,title,tag="scinex"):
    box(s,0,0,13.333,0.16,fill=CORAL,rad=False)
    txt(s,0.6,0.42,9.5,0.7,[{"runs":[(title,26,True,INK)]}])
    txt(s,10.7,0.45,2.1,0.5,[{"runs":[(tag,14,True,TEAL)]}],align=PP_ALIGN.RIGHT)
    box(s,0.6,1.18,12.13,0.02,fill=GRID,rad=False)

# ---------- Slide 1: title ----------
s=slide()
box(s,0,0,13.333,7.5,fill=WHITE,rad=False)
box(s,0,0,0.35,7.5,fill=CORAL,rad=False)
box(s,0,7.1,13.333,0.4,fill=INK,rad=False)
txt(s,1.1,2.25,11,1.4,[{"runs":[("scinex",66,True,INK)]}])
txt(s,1.13,3.62,11,0.7,[{"runs":[("PDF \u2192 структурированное химическое знание",28,True,TEAL)]}])
txt(s,1.15,4.45,10.8,1.0,[
  {"runs":[("Мульти-агентный движок научной литературы",16,False,MUTED)],"space":4},
  {"runs":[("grounded-извлечение  ·  верифицируемый граф с провенансом  ·  self-hosted",16,False,MUTED)]},
])
txt(s,1.15,7.16,11,0.3,[{"runs":[("Datacon 2026  ·  Chemistry + AI",12,True,WHITE)]}])
# ---------- Slide 2: system / pipeline ----------
s=slide(); box(s,0,0,13.333,7.5,fill=WHITE,rad=False)
header(s,"\u0421\u0438\u0441\u0442\u0435\u043c\u0430: \u043e\u0442 \u0441\u044b\u0440\u043e\u0433\u043e PDF \u043a \u0432\u0435\u0440\u0438\u0444\u0438\u0446\u0438\u0440\u0443\u0435\u043c\u043e\u043c\u0443 \u0433\u0440\u0430\u0444\u0443")
txt(s,0.6,1.35,12.1,0.5,[{"runs":[("\u041a\u0430\u0436\u0434\u044b\u0439 \u0444\u0430\u043a\u0442 \u2014 \u0441 \u0446\u0438\u0442\u0430\u0442\u043e\u0439-\u043f\u0440\u043e\u0432\u0435\u043d\u0430\u043d\u0441\u043e\u043c, \u0432\u0430\u043b\u0438\u0434\u0438\u0440\u0443\u0435\u0442\u0441\u044f \u0434\u0435\u0442\u0435\u0440\u043c\u0438\u043d\u0438\u0440\u043e\u0432\u0430\u043d\u043d\u044b\u043c \u0430\u0440\u0431\u0438\u0442\u0440\u043e\u043c, \u0441\u0445\u043b\u043e\u043f\u044b\u0432\u0430\u0435\u0442\u0441\u044f \u0432 \u043a\u0430\u043d\u043e\u043d\u0438\u0447\u0435\u0441\u043a\u0438\u0439 \u0443\u0437\u0435\u043b.",14,False,MUTED)]}])
stages=[("PDF","\u0441\u0442\u0430\u0442\u044c\u044f",MUTED),("OCR-\u0434\u0438\u0441\u043f\u0435\u0442\u0447\u0435\u0440","PyMuPDF\u00b7Mathpix\u00b7Mistral",TEAL),
        ("\u0418\u0437\u0432\u043b\u0435\u0447\u0435\u043d\u0438\u0435","LLM \u00b7 grounded",TEAL),("\u0412\u0430\u043b\u0438\u0434\u0430\u0446\u0438\u044f","RDKit \u00b7 OPSIN",TEAL),
        ("\u0413\u0440\u0430\u0444-\u0441\u0443\u0431\u0441\u0442\u0440\u0430\u0442","\u0443\u0437\u043b\u044b\u00b7\u0440\u0451\u0431\u0440\u0430\u00b7\u0440\u043e\u043b\u0438",TEAL),("\u041f\u0440\u043e\u0435\u043a\u0446\u0438\u0438","\u043c\u043e\u043b\u0435\u043a\u0443\u043b\u044b\u00b7\u0434\u0430\u043d\u043d\u044b\u0435\u00b7\u0434\u043e\u043a",CORAL)]
x=0.5; bw=1.85; gap=0.176; y=2.15; bh=1.05
for i,(nm,sub,col) in enumerate(stages):
    box(s,x,y,bw,bh,fill=col)
    txt(s,x,y+0.2,bw,0.4,[{"runs":[(nm,12.5,True,WHITE)]}],align=PP_ALIGN.CENTER)
    txt(s,x,y+0.6,bw,0.4,[{"runs":[(sub,8.5,False,WHITE)]}],align=PP_ALIGN.CENTER)
    if i<len(stages)-1:
        txt(s,x+bw-0.03,y+0.34,gap+0.06,0.4,[{"runs":[("\u2192",16,True,MUTED)]}],align=PP_ALIGN.CENTER)
    x+=bw+gap
chips=["\u041f\u0440\u043e\u0432\u0435\u043d\u0430\u043d\u0441 \u043d\u0430 \u043a\u0430\u0436\u0434\u043e\u043c \u0444\u0430\u043a\u0442\u0435","Confidence-gated \u043c\u0430\u0440\u0448\u0440\u0443\u0442\u0438\u0437\u0430\u0446\u0438\u044f","RDKit/InChIKey-\u0434\u0435\u0434\u0443\u043f","Self-hosted \u00b7 offline-ready"]
cx=0.5; cy=3.75
for c in chips:
    w=0.3+len(c)*0.094
    box(s,cx,cy,w,0.45,fill=PANEL); txt(s,cx,cy+0.11,w,0.3,[{"runs":[(c,11,True,INK)]}],align=PP_ALIGN.CENTER)
    cx+=w+0.22
txt(s,0.6,4.8,12.1,0.4,[{"runs":[("\u0422\u0440\u0438 \u043f\u0440\u043e\u0435\u043a\u0446\u0438\u0438 \u043e\u0434\u043d\u043e\u0433\u043e \u0441\u0443\u0431\u0441\u0442\u0440\u0430\u0442\u0430:",13,True,INK)]}])
projs=[("\u043c\u043e\u043b\u0435\u043a\u0443\u043b\u044b","\u0434\u0435\u0434\u0443\u043f \u043f\u043e InChIKey + SMILES + \u0432\u0430\u043b\u0438\u0434\u0430\u0446\u0438\u044f \u2014 \u0438\u0441\u0441\u043b\u0435\u0434\u043e\u0432\u0430\u043d\u0438\u0435 / substructure"),
       ("\u0434\u0430\u0442\u0430\u0441\u0435\u0442\u044b","\u043f\u043b\u043e\u0441\u043a\u0438\u0435 \u0441\u0442\u0440\u043e\u043a\u0438: metric + unit + analyte + role + conditions"),
       ("\u0434\u043e\u043a\u0443\u043c\u0435\u043d\u0442","\u0441\u0435\u043a\u0446\u0438\u0438 + \u0430\u0441\u0441\u0435\u0442\u044b (entities, data_relevance) + \u0446\u0438\u0442\u0430\u0442\u044b")]
py=5.25
for nm,d in projs:
    box(s,0.6,py,0.12,0.42,fill=CORAL,rad=False)
    txt(s,0.85,py+0.03,3.0,0.4,[{"runs":[(nm,13,True,TEAL)]}])
    txt(s,3.4,py+0.05,9.3,0.4,[{"runs":[(d,12,False,MUTED)]}])
    py+=0.55
# ---------- Slide 3: extracted structures ----------
meta=json.load(open("/tmp/pptx_assets/meta.json"))
s=slide(); box(s,0,0,13.333,7.5,fill=WHITE,rad=False)
header(s,"\u0418\u0437\u0432\u043b\u0435\u0447\u0451\u043d\u043d\u044b\u0435 \u0438 \u0432\u0435\u0440\u0438\u0444\u0438\u0446\u0438\u0440\u043e\u0432\u0430\u043d\u043d\u044b\u0435 \u0441\u0442\u0440\u0443\u043a\u0442\u0443\u0440\u044b")
txt(s,0.6,1.32,12.1,0.4,[{"runs":[("SMILES \u2192 RDKit-\u043a\u0430\u043d\u043e\u043d \u2192 InChIKey; \u0438\u0434\u0435\u043d\u0442\u0438\u0447\u043d\u043e\u0441\u0442\u044c \u0441\u043e\u0432\u043f\u0430\u0434\u0430\u0435\u0442 \u0441 NIST WebBook.",13,False,MUTED)]}])
cw=3.75; ch=2.5; xs=[0.6,4.79,8.98]; ys=[1.85,4.45]
for i,m in enumerate(meta):
    cx=xs[i%3]; cy=ys[i//3]; box(s,cx,cy,cw,ch,fill=PANEL)
    iw=2.0; ih=iw*340/460
    s.shapes.add_picture(m["png"], Inches(cx+(cw-iw)/2), Inches(cy+0.12), width=Inches(iw), height=Inches(ih))
    txt(s,cx+0.1,cy+ih+0.16,cw-0.2,0.32,[{"runs":[(m["name"],14,True,INK)]}],align=PP_ALIGN.CENTER)
    txt(s,cx+0.1,cy+ih+0.5,cw-0.2,0.28,[{"runs":[(m["ik"],10.5,False,TEAL,MONO)]}],align=PP_ALIGN.CENTER)
    tag=m["src"]; tc=CORAL if "paper" in tag else TEAL
    txt(s,cx+0.1,cy+ih+0.76,cw-0.2,0.26,[{"runs":[("\u2713 "+tag,9.5,True,tc)]}],align=PP_ALIGN.CENTER)
# ---------- Slide 4: structured output ----------
s=slide(); box(s,0,0,13.333,7.5,fill=WHITE,rad=False)
header(s,"\u0421\u0442\u0440\u0443\u043a\u0442\u0443\u0440\u0438\u0440\u043e\u0432\u0430\u043d\u043d\u044b\u0439 \u0432\u044b\u0432\u043e\u0434 \u2014 \u0444\u043e\u0440\u043c\u0430\u0442 \u043a\u0430\u043a \u0432 \u043b\u0435\u043a\u0446\u0438\u0438")
box(s,0.6,1.5,6.0,4.55,fill=CODEBG)
txt(s,0.85,1.68,5.5,0.35,[{"runs":[("figure-JSON  (\u0441\u043b\u0430\u0439\u0434 2)",12,True,CORAL)]}])
fig_json='{\n  "figure": "Figure 3",\n  "page": 4,\n  "caption": "...characterization of\n       PA membrane",\n  "image": "p4n3.png",\n  "entities": [\n     "polyamide membrane",\n     "single-walled carbon nanotubes"\n  ],\n  "data_relevance": "performance_plot"\n}'
txt(s,0.85,2.12,5.6,3.8,[{"runs":[(fig_json,11.5,False,CODEFG,MONO)]}])
box(s,6.9,1.5,5.85,2.1,fill=CODEBG)
txt(s,7.15,1.68,5.4,0.35,[{"runs":[("\u0441\u0442\u0440\u043e\u043a\u0430 \u0434\u0430\u0442\u0430\u0441\u0435\u0442\u0430  (mention\u2192role)",12,True,CORAL)]}])
ds='{ "metric": "yield", "value": 78.0,\n  "unit": "%",\n  "entity": "2-amino-4H-chromene",\n  "inchikey": "ARNCZJZLEMLOBH-...",\n  "role": "product" }'
txt(s,7.15,2.1,5.5,1.45,[{"runs":[(ds,11.5,False,CODEFG,MONO)]}])
box(s,6.9,3.8,5.85,2.25,fill=PANEL)
txt(s,7.15,3.98,5.4,0.35,[{"runs":[("\u0433\u0440\u0430\u0444-\u0441\u0443\u0431\u0441\u0442\u0440\u0430\u0442 (1 \u0441\u0442\u0430\u0442\u044c\u044f)",12,True,INK)]}])
stats=[("\u0441\u0443\u0449\u043d\u043e\u0441\u0442\u0438","11"),("\u043c\u043e\u043b\u0435\u043a\u0443\u043b\u044b","7"),("\u0438\u0437\u043c\u0435\u0440.","6"),("\u0441\u0432\u044f\u0437\u0438","5"),("\u0444\u0438\u0433\u0443\u0440\u044b","5")]
sx=7.15
for nm,v in stats:
    box(s,sx,4.45,1.0,1.05,fill=WHITE)
    txt(s,sx,4.57,1.0,0.5,[{"runs":[(v,20,True,TEAL)]}],align=PP_ALIGN.CENTER)
    txt(s,sx-0.05,5.12,1.1,0.35,[{"runs":[(nm,9,False,MUTED)]}],align=PP_ALIGN.CENTER)
    sx+=1.12
txt(s,6.9,6.2,5.85,0.5,[{"runs":[("\u043a\u0430\u0436\u0434\u044b\u0439 \u0443\u0437\u0435\u043b/\u0440\u0435\u0431\u0440\u043e: \u043f\u0440\u043e\u0432\u0435\u043d\u0430\u043d\u0441 (paper+span / figure+bbox) + validation",10,False,MUTED)]}])
txt(s,0.6,6.2,6.0,0.5,[{"runs":[("\u0440\u0435\u0430\u043b\u044c\u043d\u044b\u0439 \u0432\u044b\u0432\u043e\u0434 \u043f\u0430\u0439\u043f\u043b\u0430\u0439\u043d\u0430 (Nature Comms, solar-lithium)",10,False,MUTED)]}])
# ---------- Slide 5: NIST mini-task -> dataset row ----------
s=slide(); box(s,0,0,13.333,7.5,fill=WHITE,rad=False)
header(s,"Mini-task: NIST WebBook → строка датасета")
txt(s,0.6,1.28,12.2,0.4,[{"runs":[("HTML-страница NIST WebBook → извлечение полей → одна валидная строка (мини-задание, A–F)",12,False,MUTED)]}])
NROWS=[("Aspirin","C9H8O4","180.1574","BSYNRYMUTXBXSQ-UHFFFAOYSA-N","50-78-2"),
       ("Caffeine","C8H10N4O2","194.1906","RYYVLZVUVIJVGH-UHFFFAOYSA-N","58-08-2"),
       ("Ibuprofen","C13H18O2","206.2808","HEFNNWSXXWATRW-UHFFFAOYSA-N","15687-27-1"),
       ("Benzene","C6H6","78.1118","UHOVQNZJYSORNB-UHFFFAOYSA-N","71-43-2"),
       ("Acetone","C3H6O","58.0791","CSCPPACGZOOCGX-UHFFFAOYSA-N","67-64-1"),
       ("Ethanol","C2H6O","46.0684","LFQSCWFLJHTTHZ-UHFFFAOYSA-N","64-17-5")]
_hdr=f"{'name':<11}{'formula':<11}{'mw':<11}{'inchikey':<30}{'cas_number':<12}{'source':<6}"
box(s,0.6,1.72,12.13,2.5,fill=CODEBG)
_par=[{"runs":[(_hdr,11,True,CODEKEY,MONO)],"space":5}]
for _n,_f,_m,_ik,_c in NROWS:
    _par.append({"runs":[(f"{_n:<11}{_f:<11}{_m:<11}{_ik:<30}{_c:<12}{'NIST':<6}",11,False,CODEFG,MONO)],"space":2})
txt(s,0.85,1.9,11.8,2.2,_par)
box(s,0.6,4.45,12.13,2.28,fill=CODEBG)
txt(s,0.85,4.57,8.0,0.3,[{"runs":[("полная строка (все 7 полей) — Aspirin",12,True,CORAL)]}])
_rec=[("name","Aspirin"),("formula","C9H8O4"),("mw","180.1574"),
      ("inchi","InChI=1S/C9H8O4/c1-6(10)13-8-5-3-2-4-7(8)9(11)12/h2-5H,1H3,(H,11,12)"),
      ("inchikey","BSYNRYMUTXBXSQ-UHFFFAOYSA-N"),("cas_number","50-78-2"),
      ("source_url","https://webbook.nist.gov/cgi/cbook.cgi?Name=Aspirin&Units=SI")]
_rp=[]
for _k,_v in _rec:
    _rp.append({"runs":[(f"{_k:<12}",11,False,CODEKEY,MONO),(_v,11,False,CODEFG,MONO)],"space":2})
txt(s,0.85,4.98,11.8,1.7,_rp)
txt(s,0.6,6.86,12.2,0.4,[{"runs":[("InChIKey из NIST == RDKit-канон (кросс-валидация источника)  ·  enrich/nist.py",10,False,MUTED)]}])

# ---------- Slide 6: OCSR (image -> SMILES), beyond the lecture ----------
meta=json.load(open("/tmp/pptx_assets/meta.json"))
asp=[m for m in meta if m["ik"]=="BSYNRYMUTXBXSQ"][0]
s=slide(); box(s,0,0,13.333,7.5,fill=WHITE,rad=False)
header(s,"OCSR: структура с рисунка → SMILES   (за рамками лекции)")
txt(s,0.6,1.32,12.1,0.4,[{"runs":[("Два независимых ридера читают кроп фигуры; консенсус по InChIKey — арбитр RDKit.",13,False,MUTED)]}])
box(s,0.6,2.0,2.7,2.5,fill=PANEL,line=GRID)
s.shapes.add_picture(asp["png"], Inches(0.95), Inches(2.2), width=Inches(2.0), height=Inches(2.0*340/460))
txt(s,0.6,4.18,2.7,0.3,[{"runs":[("figure crop (PDF)",10,True,MUTED)]}],align=PP_ALIGN.CENTER)
txt(s,3.35,3.0,0.5,0.4,[{"runs":[("→",20,True,MUTED)]}],align=PP_ALIGN.CENTER)
box(s,3.95,2.0,4.2,1.15,fill=CODEBG)
txt(s,4.15,2.14,3.8,0.3,[{"runs":[("vision (Opus + Gemini)",11,True,CORAL)]}])
txt(s,4.15,2.52,3.9,0.5,[{"runs":[("CC(=O)Oc1ccccc1C(=O)O",12,False,CODEFG,MONO)]}])
box(s,3.95,3.35,4.2,1.15,fill=CODEBG)
txt(s,4.15,3.49,3.8,0.3,[{"runs":[("DECIMER (CNN, off-line)",11,True,CORAL)]}])
txt(s,4.15,3.87,3.9,0.5,[{"runs":[("CC(=O)Oc1ccccc1C(=O)O",12,False,CODEFG,MONO)]}])
txt(s,8.2,3.0,0.5,0.4,[{"runs":[("→",20,True,MUTED)]}],align=PP_ALIGN.CENTER)
box(s,8.75,2.0,3.98,2.5,fill=TEAL)
txt(s,8.95,2.22,3.6,0.35,[{"runs":[("✓ консенсус (agree)",14,True,WHITE)]}])
txt(s,8.95,2.8,3.7,1.5,[
  {"runs":[("RDKit-канон → InChIKey",10,False,WHITE)],"space":6},
  {"runs":[("BSYNRYMUTXBXSQ",13,True,WHITE,MONO)],"space":6},
  {"runs":[("status: ok    conf: 0.95",11,False,WHITE,MONO)]},
])
box(s,0.6,4.95,12.13,1.3,fill=PANEL)
txt(s,0.85,5.12,11.6,1.05,[
  {"runs":[("Почему консенсус, а не один ридер: ",12,True,INK),("ошибки DECIMER некоррелированы с vision → совпадение ловит «valid-but-wrong» (RDKit-валидно, но не та молекула).",12,False,MUTED)],"space":7},
  {"runs":[("Робастность: ",12,True,INK),("DECIMER устойчив к шуму (4/4 зашумлённых); frontier-VLM теряют сложные каркасы под шумом (урок camphor).",12,False,MUTED)]},
])

# ---------- Slides 7-9: selector graph, one stage per slide ----------
def _graph_slide(png, title, subtitle, legend=False):
    if not png: return
    from PIL import Image as _IMG
    iw,ih=_IMG.open(png).size
    s=slide(); box(s,0,0,13.333,7.5,fill=WHITE,rad=False)
    header(s,title)
    if subtitle: txt(s,0.6,1.28,12.2,0.4,[{"runs":[(subtitle,12,False,MUTED)]}])
    aw,ah=12.33,(4.7 if legend else 5.1)
    gw=min(aw, ah*iw/ih); gh=gw*ih/iw
    gx=0.5+(aw-gw)/2; gy=1.95+(ah-gh)/2
    s.shapes.add_picture(png, Inches(gx), Inches(gy), width=Inches(gw))
    if legend:
        ly=6.62
        box(s,3.0,ly,0.3,0.3,fill=RGBColor(0xE8,0xF5,0xE9),line=RGBColor(0x2E,0x7D,0x32))
        txt(s,3.4,ly+0.03,3.6,0.3,[{"runs":[("\u0434\u0435\u0442\u0435\u0440\u043c\u0438\u043d\u0438\u0440\u043e\u0432\u0430\u043d\u043d\u044b\u0439 \u0430\u0440\u0431\u0438\u0442\u0440 (RDKit/OPSIN)",10,False,INK)]}])
        box(s,7.4,ly,0.3,0.3,fill=RGBColor(0xFF,0xF3,0xE0),line=RGBColor(0xE6,0x51,0x00))
        txt(s,7.8,ly+0.03,2.4,0.3,[{"runs":[("\u043e\u0447\u0435\u0440\u0435\u0434\u044c \u0440\u0435\u0432\u044c\u044e",10,False,INK)]}])

_graph_slide(SG1_PNG, "\u0413\u0440\u0430\u0444 \u2014 \u0441\u0442\u0430\u0434\u0438\u044f 1: \u043f\u0440\u0438\u0451\u043c \u0438 OCR",
             "PyMuPDF $0 \u0434\u043b\u044f \u0446\u0438\u0444\u0440\u043e\u0432\u044b\u0445 \u00b7 Mistral \u0434\u043b\u044f \u0441\u043a\u0430\u043d\u043e\u0432 \u00b7 Mathpix-\u044d\u0441\u043a\u0430\u043b\u0430\u0446\u0438\u044f \u043f\u043e \u0444\u043e\u0440\u043c\u0443\u043b\u0430\u043c/\u0443\u0432\u0435\u0440\u0435\u043d\u043d\u043e\u0441\u0442\u0438", legend=True)
_graph_slide(SG2_PNG, "\u0413\u0440\u0430\u0444 \u2014 \u0441\u0442\u0430\u0434\u0438\u044f 2: \u0441\u0442\u0440\u0443\u043a\u0442\u0443\u0440\u0430 \u2192 SMILES",
             "OCSR (vision \u2225 DECIMER, \u043a\u043e\u043d\u0441\u0435\u043d\u0441\u0443\u0441 \u043f\u043e InChIKey) \u0438 OPSIN (\u0438\u043c\u044f\u2192SMILES); \u0430\u0440\u0431\u0438\u0442\u0440 \u2014 RDKit", legend=True)
_graph_slide(SG3_PNG, "\u0413\u0440\u0430\u0444 \u2014 \u0441\u0442\u0430\u0434\u0438\u044f 3: QA",
             "\u0440\u0435\u0442\u0440\u0438\u0432 top-k \u043f\u043e \u0442\u0435\u043a\u0441\u0442\u0443+\u0440\u0438\u0441\u0443\u043d\u043a\u0430\u043c \u00b7 \u044d\u0441\u043a\u0430\u043b\u0430\u0446\u0438\u044f L1 \u2192 L2/L3 \u043a\u043e\u043d\u0441\u0435\u043d\u0441\u0443\u0441 \u00b7 \u0441\u0442\u0440\u0443\u043a\u0442\u0443\u0440\u0438\u0440\u043e\u0432\u0430\u043d\u043d\u044b\u0439 BenchAnswer", legend=False)

prs.save("/root/scinex/scinex_overview.pptx")
print("DECK saved:", len(prs.slides._sldIdLst), "slides")
